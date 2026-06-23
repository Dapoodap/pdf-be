import io
import os
import tempfile
import subprocess
from typing import List
from pdf2image import convert_from_bytes
from pdf2docx import Converter as Pdf2DocxConverter
import pdfplumber
import pandas as pd
from pptx import Presentation
from PIL import Image

# PDF to X
def pdf_to_images(file_content: bytes) -> List[bytes]:
    # Returns a list of image bytes (PNG)
    images = convert_from_bytes(file_content, fmt="png")
    image_bytes_list = []
    for img in images:
        img_byte_arr = io.BytesIO()
        img.save(img_byte_arr, format='PNG')
        image_bytes_list.append(img_byte_arr.getvalue())
    return image_bytes_list

def pdf_to_docx(file_content: bytes) -> bytes:
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_pdf:
        temp_pdf.write(file_content)
        temp_pdf_path = temp_pdf.name
        
    temp_docx_path = temp_pdf_path.replace(".pdf", ".docx")
    
    try:
        cv = Pdf2DocxConverter(temp_pdf_path)
        cv.convert(temp_docx_path)
        cv.close()
        
        with open(temp_docx_path, "rb") as f:
            docx_bytes = f.read()
            
        return docx_bytes
    finally:
        if os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)
        if os.path.exists(temp_docx_path):
            os.remove(temp_docx_path)

def pdf_to_xlsx(file_content: bytes) -> bytes:
    # Extracts tables and saves to a single excel file with multiple sheets (one per page)
    output = io.BytesIO()
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_pdf:
        temp_pdf.write(file_content)
        temp_pdf_path = temp_pdf.name
        
    try:
        with pdfplumber.open(temp_pdf_path) as pdf:
            with pd.ExcelWriter(output, engine='openpyxl') as writer:
                for i, page in enumerate(pdf.pages):
                    table = page.extract_table()
                    if table:
                        # Assuming the first row is the header if table exists
                        if len(table) > 1:
                            df = pd.DataFrame(table[1:], columns=table[0])
                        else:
                            df = pd.DataFrame(table)
                        df.to_excel(writer, sheet_name=f"Page {i+1}", index=False)
                    else:
                        # Write empty dataframe if no table found to keep page indexing
                        df = pd.DataFrame(["No table found on this page"])
                        df.to_excel(writer, sheet_name=f"Page {i+1}", index=False, header=False)
        return output.getvalue()
    finally:
        if os.path.exists(temp_pdf_path):
            os.remove(temp_pdf_path)

def pdf_to_pptx(file_content: bytes) -> bytes:
    # Convert PDF pages to images and embed in PPTX
    images = convert_from_bytes(file_content, fmt="png")
    prs = Presentation()
    
    blank_slide_layout = prs.slide_layouts[6]
    
    for img in images:
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as temp_img:
            img.save(temp_img.name, format="PNG")
            temp_img_path = temp_img.name
            
        slide = prs.slides.add_slide(blank_slide_layout)
        
        width = prs.slide_width
        height = prs.slide_height
        
        slide.shapes.add_picture(temp_img_path, 0, 0, width, height)
        os.remove(temp_img_path)
        
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()

# X to PDF
def x_to_pdf(file_content: bytes, file_ext: str) -> bytes:
    if file_ext.lower() in [".png", ".jpg", ".jpeg"]:
        img = Image.open(io.BytesIO(file_content))
        if img.mode != 'RGB':
            img = img.convert('RGB')
        output = io.BytesIO()
        img.save(output, format="PDF")
        return output.getvalue()
        
    # For DOCX, XLSX, PPTX using LibreOffice Headless
    with tempfile.NamedTemporaryFile(suffix=file_ext, delete=False) as temp_input:
        temp_input.write(file_content)
        temp_input_path = temp_input.name
        
    temp_dir = tempfile.gettempdir()
    
    try:
        process = subprocess.run(
            [
                "libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                temp_input_path,
                "--outdir",
                temp_dir,
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=120,  # Kill process jika lebih dari 2 menit (mencegah zombie process)
        )

        if process.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {process.stderr.decode()}")
            
        base_name = os.path.splitext(os.path.basename(temp_input_path))[0]
        output_pdf_path = os.path.join(temp_dir, f"{base_name}.pdf")
        
        if not os.path.exists(output_pdf_path):
            raise RuntimeError("LibreOffice did not produce a PDF file.")
            
        with open(output_pdf_path, "rb") as f:
            pdf_bytes = f.read()
            
        os.remove(output_pdf_path)
        return pdf_bytes
    except subprocess.TimeoutExpired:
        raise RuntimeError("LibreOffice conversion timed out after 120 seconds. Please try a smaller file.")
    finally:
        if os.path.exists(temp_input_path):
            os.remove(temp_input_path)
